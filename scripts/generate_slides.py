"""
scripts/generate_slides.py
============================
Generates presentation.pptx — the presentation slide deck for the Smart
Factories course (Tor Vergata, Prof. Annalisa Santolamazza), combining:
  Track 2 — Digital Twin, Data-driven (Quality Monitoring)
  Track 3 — RAG for Industrial Knowledge Management (Quality Assistant)

Run from the project root (after scripts/train.py and scripts/make_assets.py):
    python scripts/generate_slides.py

Requires: python-pptx
"""

import io
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
OUT_PPTX = ROOT / "presentation.pptx"

with open(ROOT / "models" / "cv_results.json") as f:
    RES = json.load(f)

# ── Colours ──────────────────────────────────────────────────────────────
GREEN  = RGBColor(0x1E, 0x71, 0x45)
DGREEN = RGBColor(0x0B, 0x4A, 0x2A)
LGREEN = RGBColor(0xE3, 0xF1, 0xE8)
TEAL   = RGBColor(0x00, 0x69, 0x5C)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
DGRAY  = RGBColor(0x42, 0x42, 0x42)
MGRAY  = RGBColor(0x75, 0x75, 0x75)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = color
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, size=18, bold=False,
              color=DGRAY, align=PP_ALIGN.LEFT, font="Calibri", italic=False,
              anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, size=15, color=DGRAY,
                 space_after=8, bullet_color=GREEN):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = level
        para.space_after = Pt(space_after)
        prefix = "▪ " if level == 0 else "– "
        run = para.add_run()
        run.text = prefix + text
        run.font.size = Pt(size - level * 1.5)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_image(slide, path, left, top, max_w, max_h):
    with PILImage.open(path) as im:
        w, h = im.size
    ratio = min(max_w / w, max_h / h)
    nw, nh = int(w * ratio), int(h * ratio)
    pic_left = left + int((max_w - nw) / 2)
    pic_top = top + int((max_h - nh) / 2)
    slide.shapes.add_picture(str(path), pic_left, pic_top, width=nw, height=nh)


def add_image_downscaled(slide, path, left, top, max_w, max_h, max_px=900):
    """Like add_image, but re-encodes large source photos so the pptx stays small."""
    with PILImage.open(path) as im:
        im = im.convert("RGB")
        w, h = im.size
        if max(w, h) > max_px:
            scale = max_px / max(w, h)
            im = im.resize((int(w * scale), int(h * scale)))
            w, h = im.size
        buf = io.BytesIO()
        im.save(buf, format="JPEG", quality=85)
        buf.seek(0)
    ratio = min(max_w / w, max_h / h)
    nw, nh = int(w * ratio), int(h * ratio)
    pic_left = left + int((max_w - nw) / 2)
    pic_top = top + int((max_h - nh) / 2)
    slide.shapes.add_picture(buf, pic_left, pic_top, width=nw, height=nh)


def add_header(slide, title, track=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), GREEN if track != "rag" else TEAL)
    add_text(slide, Inches(0.5), Inches(0.12), Inches(11), Inches(0.9),
             title, size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, Inches(7.35), SLIDE_W, Inches(0.15), GREEN if track != "rag" else TEAL)


def add_footer(slide, page_num):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             "Checkered Fabric AOI Digital Twin — Smart Factories A.Y. 2025-2026",
             size=9, color=MGRAY, italic=True)
    add_text(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.3),
             str(page_num), size=9, color=MGRAY, align=PP_ALIGN.RIGHT)


def content_slide(title, track=None):
    s = new_slide()
    add_header(s, title, track)
    return s


def track_band(slide, text, color):
    band = add_rect(slide, Inches(0.5), Inches(1.35), Inches(12.33), Inches(0.5), color)
    add_text(slide, Inches(0.5), Inches(1.35), Inches(12.33), Inches(0.5), text,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# 15-query RAG evaluation summary (matches scripts/eval_rag.py)
EVAL_SUMMARY = [
    ("Correct", 12, "80%"),
    ("Partial", 2, "13%"),
    ("No answer", 1, "7%"),
    ("Hallucinated", 0, "0%"),
]

CORPUS_DOCS = [
    "Process Overview — Checkered Fabric Weaving and Inspection Line",
    "ISO 23247 Mapping — Digital Twin Reference Architecture",
    "What Makes This a Data-Driven Digital Twin",
    "Dataset and Model — 150 Images, ResNet18, Grad-CAM",
    "Model Performance — 5-Fold Cross-Validation Results",
    "SOP — Circle Defect (Hole / Puncture)",
    "SOP — Line Defect (Scratch / Cut Mark)",
    "SOP — “No Defect” Classification and False-Negative Risk",
    "Quality Standard Reference — 4-Point Fabric Inspection System",
    "AOI Station Maintenance and Calibration Guide",
    "Deployment and Integration in a Real Industrial System",
    "Critical Discussion — Dataset and Model Limitations",
]


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DGREEN)
add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(0.06), RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(1), Inches(1.6), Inches(11.33), Inches(1.2),
         "Checkered Fabric AOI", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.5), Inches(11.33), Inches(1.0),
         "Digital Twin & Quality Assistant", size=30, bold=True, color=RGBColor(0xB8, 0xE0, 0xC8),
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.6), Inches(11.33), Inches(0.5),
         "Track 2 — Digital Twin, Data-Driven (Quality Monitoring)  +  "
         "Track 3 — RAG for Industrial Knowledge Management",
         size=15, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(1), Inches(5.8), Inches(11.33), Inches(0.5),
         "Kuhyar Saeedi  ·  Danial Mahmoody  ·  Davood Jokar  ·  Mahyar Emami  ·  Nima Shahrokhi",
         size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.4), Inches(11.33), Inches(0.5),
         "Smart Factories — A.Y. 2025-2026  |  Prof. Annalisa Santolamazza  |  "
         "Universita degli Studi di Roma Tor Vergata",
         size=13, color=RGBColor(0xCC, 0xDD, 0xD3), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CASE INTRODUCTION
# ════════════════════════════════════════════════════════════════════════
s = content_slide("The Case: An AOI Station on a Checkered Fabric Line")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.6),
         "Yarn is woven into a checkered (gingham) cotton fabric. A camera at the "
         "finishing/inspection frame captures images of the moving fabric surface.",
         size=16, color=DGRAY)
add_bullets(s, Inches(0.6), Inches(2.3), Inches(6.0), Inches(4.5), [
    "Each image patch must be classified into one of three categories:",
    ("Circle — a roughly circular hole / puncture in the weave", 1),
    ("Line — a linear scratch, cut, or abrasion mark", 1),
    ("No defect — a clean, regular checkered pattern", 1),
    "Two complementary tracks, one shared dashboard:",
    ("Track 2: train + evaluate a defect classifier (Digital Twin Entity)", 1),
    ("Track 3: Quality Assistant answers “what caused this, what do I do?”", 1),
], size=16)
# right-side class examples
img_files = {"Circle": "1.jpg", "Line": "1.jpg", "No defect": "1.jpg"}
classes_dir = ROOT / "dataset"
x = Inches(7.0)
for i, cls in enumerate(["Circle", "Line", "No defect"]):
    folder = classes_dir / cls
    samples = sorted(folder.glob("*.jpg"))
    if samples:
        add_image_downscaled(s, samples[0], x + Inches(i * 2.1), Inches(2.3), Inches(2.0), Inches(2.0))
    add_text(s, x + Inches(i * 2.1), Inches(4.4), Inches(2.0), Inches(0.4), cls,
             size=14, bold=True, color=DGRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(7.0), Inches(5.0), Inches(6.2), Inches(1.8),
         "The Defect Classifier (Track 2) and Quality Assistant (Track 3) are integrated "
         "into a single 5-page Streamlit dashboard: a defect prediction automatically "
         "surfaces the relevant SOP from the knowledge base.",
         size=13, color=DGRAY, italic=True)
add_footer(s, 2)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ISO 23247 MAPPING
# ════════════════════════════════════════════════════════════════════════
s = content_slide("ISO 23247 Reference Architecture Mapping")
add_text(s, Inches(0.5), Inches(1.35), Inches(12.33), Inches(0.6),
         "ISO 23247 defines five entity types for a manufacturing digital twin. "
         "Mapping the project onto them is what makes it a digital twin, not just an ML model.",
         size=15, color=DGRAY)

iso_rows = [
    ("Observable Manufacturing\nElement (OME)", "The fabric web at the inspection frame"),
    ("Device Communication\nEntity (DCE)", "AOI camera / frame grabber capturing fabric images"),
    ("Data Collection &\nDevice Control", "Pre-processing: resize to 224x224, normalisation, augmentation"),
    ("Digital Twin Entity\n(Core)", "ResNet18 classifier + Grad-CAM explainability module"),
    ("User Entity", "Streamlit dashboard — Defect Classifier + Quality Assistant pages"),
    ("Cross System Entity", "MES / QMS for traceability (discussed conceptually for deployment)"),
]
tbl_top = Inches(2.1)
tbl_h = Inches(4.9)
rows, cols = len(iso_rows) + 1, 2
gtbl = s.shapes.add_table(rows, cols, Inches(0.7), tbl_top, Inches(11.9), tbl_h).table
gtbl.columns[0].width = Inches(3.5)
gtbl.columns[1].width = Inches(8.4)
hdr_cells = ["ISO 23247 Entity", "This Project"]
for c, txt in enumerate(hdr_cells):
    cell = gtbl.cell(0, c)
    cell.text = txt
    cell.fill.solid()
    cell.fill.fore_color.rgb = GREEN
    para = cell.text_frame.paragraphs[0]
    para.runs[0].font.bold = True
    para.runs[0].font.color.rgb = WHITE
    para.runs[0].font.size = Pt(15)
for r, (a, b) in enumerate(iso_rows, start=1):
    for c, txt in enumerate([a, b]):
        cell = gtbl.cell(r, c)
        cell.text = txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = LGRAY if r % 2 == 0 else WHITE
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.color.rgb = DGRAY
                run.font.bold = (c == 0)
add_footer(s, 3)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TRACK 2 SECTION HEADER
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, GREEN)
add_text(s, Inches(1), Inches(2.8), Inches(11.33), Inches(1.0),
         "TRACK 2", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.9), Inches(11.33), Inches(0.8),
         "Digital Twin, Data-Driven — Quality Monitoring", size=24, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.8), Inches(11.33), Inches(0.6),
         "Dataset & EDA  ·  Model & Methodology  ·  5-Fold CV Results  ·  Grad-CAM  ·  Critical Discussion",
         size=14, color=RGBColor(0xD9, 0xEF, 0xE3), align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DATASET & EDA
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Dataset and Exploratory Data Analysis")
add_bullets(s, Inches(0.6), Inches(1.4), Inches(5.6), Inches(5.0), [
    "150 images total — perfectly balanced: 50 Circle / 50 Line / 50 No defect",
    "3072 x 3072 px photographs of the same blue/yellow checkered fabric sample",
    "Circle & Line defects created by placing a physical object (punched hole, "
    "ruler/blade) on good fabric to simulate defect signatures",
    "All images valid, readable, uniquely named, same resolution and RGB mode "
    "— no missing values or corrupt files",
    "Resized to 224x224 during pre-processing, regardless of native size",
    "EDA confirms: no class shows a systematically different mean brightness "
    "— the model cannot rely on a brightness shortcut",
], size=16, space_after=14)
add_image(s, ASSETS / "eda_overview.png", Inches(6.4), Inches(1.5), Inches(6.5), Inches(5.0))
add_footer(s, 5)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODEL & METHODOLOGY
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Model and Methodology")
add_bullets(s, Inches(0.6), Inches(1.4), Inches(12.0), Inches(5.5), [
    "Architecture: ResNet18 pretrained on ImageNet",
    ("Convolutional backbone FROZEN; final layer replaced with a fresh 3-class linear head", 1),
    ("Only the head is trained — 8 epochs, Adam (lr = 1e-3), cross-entropy loss", 1),
    "Why transfer learning + frozen backbone?",
    ("With only 150 images, fine-tuning the full network would overfit badly", 1),
    ("ImageNet features (edges, textures, shapes) transfer well to fabric images", 1),
    ("The model only has to learn the final 3-class decision boundary — a much "
     "smaller, better-posed learning problem", 1),
    "Pre-processing & augmentation",
    ("All images resized to 224x224 and converted to tensors", 1),
    ("Training augmentation: random horizontal flip, random rotation (±10°), "
     "colour jitter — reduces overfitting on the small dataset", 1),
    ("Evaluation uses a deterministic resize-only transform", 1),
    "Validation protocol: stratified 5-fold cross-validation (random_state = 42)",
    ("Each fold trains on 120 images (24/class x 5), tests on 30 (10/class)", 1),
    ("Every image used for testing exactly once across the 5 folds", 1),
], size=15, space_after=6)
add_footer(s, 6)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 5-FOLD CV RESULTS
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Results: 5-Fold Cross-Validation Accuracy")
add_image(s, ASSETS / "fold_accuracy.png", Inches(0.5), Inches(1.4), Inches(7.5), Inches(5.0))

box_top = Inches(1.6)
add_rect(s, Inches(8.3), box_top, Inches(4.5), Inches(1.6), LGREEN)
add_text(s, Inches(8.5), box_top + Inches(0.15), Inches(4.1), Inches(0.5),
         "Average Accuracy", size=16, bold=True, color=DGREEN)
add_text(s, Inches(8.5), box_top + Inches(0.55), Inches(4.1), Inches(0.9),
         f"{RES['avg_accuracy']:.1%}  (±{RES['std_accuracy']:.1%})", size=34, bold=True, color=DGREEN)

add_bullets(s, Inches(8.3), Inches(3.5), Inches(4.6), Inches(3.4), [
    f"Per-fold range: {min(RES['fold_accuracies']):.0%} – {max(RES['fold_accuracies']):.0%}",
    "150 images, 8 epochs/fold",
    "With 30 test images/fold, one misclassification = ~3.3 pp shift",
    "The observed spread is expected sampling noise for this dataset "
    "size — not model instability",
    "Headline 94% should always be reported with its std. dev.",
], size=15)
add_footer(s, 7)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CONFUSION MATRIX & CLASSIFICATION REPORT
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Confusion Matrix and Classification Report")
add_image(s, ASSETS / "confusion_matrix.png", Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.0))

rep = RES["classification_report"]
classes = RES["classes"]
rows = len(classes) + 2
cols = 5
tbl = s.shapes.add_table(rows, cols, Inches(6.6), Inches(1.6), Inches(6.2), Inches(2.4)).table
for c, txt in enumerate(["Class", "Precision", "Recall", "F1", "Support"]):
    cell = tbl.cell(0, c)
    cell.text = txt
    cell.fill.solid()
    cell.fill.fore_color.rgb = GREEN
    cell.text_frame.paragraphs[0].runs[0].font.bold = True
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
for r, c_name in enumerate(classes, start=1):
    vals = [c_name, f"{rep[c_name]['precision']:.2f}", f"{rep[c_name]['recall']:.2f}",
            f"{rep[c_name]['f1-score']:.2f}", f"{int(rep[c_name]['support'])}"]
    for c, txt in enumerate(vals):
        cell = tbl.cell(r, c)
        cell.text = str(txt)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LGRAY if r % 2 == 0 else WHITE
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(13)
        run.font.color.rgb = DGRAY
        run.font.bold = (c == 0)
macro = rep["macro avg"]
vals = ["macro avg", f"{macro['precision']:.2f}", f"{macro['recall']:.2f}",
        f"{macro['f1-score']:.2f}", f"{int(macro['support'])}"]
for c, txt in enumerate(vals):
    cell = tbl.cell(rows - 1, c)
    cell.text = str(txt)
    cell.fill.solid()
    cell.fill.fore_color.rgb = LGREEN
    run = cell.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = DGREEN

add_bullets(s, Inches(6.6), Inches(4.3), Inches(6.2), Inches(2.6), [
    f"Largest confusion: {RES['confusion_matrix'][0][2]} Circle → No defect, "
    f"{RES['confusion_matrix'][2][0]} No defect → Circle",
    "Line classified perfectly (precision = recall = 1.00)",
    "A small/faint hole can visually resemble a normal weave gap under "
    "uneven lighting — a linear scratch is more visually distinctive",
], size=15)
add_footer(s, 8)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — GRAD-CAM EXPLAINABILITY
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Explainability: Grad-CAM")
add_text(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.9),
         "Grad-CAM on the last residual block (layer4) highlights the image regions "
         "most responsible for the predicted class — lets an operator verify the model "
         "is reacting to the actual defect, not lighting/staging artefacts.",
         size=15, color=DGRAY)
add_image(s, ASSETS / "gradcam_panel.png", Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.7))
add_footer(s, 9)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CRITICAL DISCUSSION TRACK 2
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Critical Discussion — Track 2")
add_bullets(s, Inches(0.6), Inches(1.4), Inches(12.0), Inches(5.6), [
    f"Dataset size: only {RES['n_images']} images (50/class). 5-fold CV gives a "
    f"robust relative comparison, but std = {RES['std_accuracy']:.1%} means the "
    "absolute accuracy estimate has wide uncertainty",
    "Staged vs. real defects: Circle/Line defects were created by placing physical "
    "objects on good fabric — real weaving defects (yarn pull-outs, abrasion) may "
    "have different edge/shadow signatures",
    "Frozen backbone caps achievable accuracy — fine-tuning deeper layers with more "
    "data could improve Circle vs. No-defect discrimination",
    "Single fabric pattern: the model has only seen this one blue/yellow check — a "
    "different weave/colour would need additional labelled data and retraining",
    "Real deployment: edge device runs the same pipeline at line speed; each "
    "prediction (class, confidence, Grad-CAM region, roll ID, timestamp) is logged "
    "to the MES/QMS for traceability and SPC charting",
    "A confidence threshold (e.g. 0.7-0.8) would route low-confidence predictions "
    "to human review rather than auto-classifying",
], size=17, space_after=14)
add_footer(s, 10)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TRACK 3 SECTION HEADER
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, TEAL)
add_text(s, Inches(1), Inches(2.8), Inches(11.33), Inches(1.0),
         "TRACK 3", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.9), Inches(11.33), Inches(0.8),
         "RAG for Industrial Knowledge Management — Quality Assistant", size=24, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.8), Inches(11.33), Inches(0.6),
         "Use Case  ·  Corpus Design  ·  System Architecture  ·  15-Query Evaluation  ·  Discussion",
         size=14, color=RGBColor(0xCC, 0xEE, 0xE8), align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 12 — USE CASE
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Use Case and User Needs", track="rag")
add_text(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.9),
         "Operators and quality engineers need fast, natural-language access to SOPs, "
         "corrective actions, and quality standards — without searching paper manuals.",
         size=16, color=DGRAY)
add_bullets(s, Inches(0.6), Inches(2.3), Inches(12.0), Inches(4.7), [
    "Operator who just received a Circle/Line prediction wants to immediately know "
    "the likely root causes and corrective action — closing the loop from detection "
    "to action",
    "New operator wants to understand how the digital twin works (ISO 23247 mapping, "
    "model performance, limitations) without reading full documentation",
    "Quality engineer wants to relate AOI output to the mill's existing point-grading "
    "quality standard (4-point fabric inspection system)",
    "Technician wants the AOI station's maintenance/calibration checklist on demand "
    "at the inspection frame",
], size=17, space_after=16)
add_footer(s, 12)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CORPUS DESIGN
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Corpus Design — 12-Document Knowledge Base", track="rag")
add_text(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.7),
         "Illustrative SOP/standards corpus written for this project. Each document "
         "is a short (200-400 word) retrievable unit — no chunking needed at this scale, "
         "which keeps citations interpretable.",
         size=15, color=DGRAY)

half = (len(CORPUS_DOCS) + 1) // 2
col1 = CORPUS_DOCS[:half]
col2 = CORPUS_DOCS[half:]
add_bullets(s, Inches(0.6), Inches(2.2), Inches(6.2), Inches(4.8),
            [f"{i+1}. {d}" for i, d in enumerate(col1)], size=14, space_after=10)
add_bullets(s, Inches(6.9), Inches(2.2), Inches(6.2), Inches(4.8),
            [f"{half+i+1}. {d}" for i, d in enumerate(col2)], size=14, space_after=10)
add_footer(s, 13)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 14 — SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════
s = content_slide("System Architecture: Retrieval + Generation", track="rag")

stages = [
    ("1. Query", "Operator asks a question\n(chat input or auto-triggered\nafter a defect prediction)"),
    ("2. Retrieval", "TF-IDF vector search over\n12 documents (cosine similarity)\n→ top-3 ranked results\n\n(optional: sentence-transformers\nsemantic embeddings)"),
    ("3. Generation", "Default: extractive — top\nsentences assembled from\nretrieved docs (offline, zero\nhallucination risk)\n\nOptional: Claude Haiku 4.5,\ngrounded by context-only prompt"),
    ("4. Answer + Sources", "Chat response with an\nexpandable “Sources” panel\nshowing retrieved docs and\nrelevance scores"),
]
box_w = Inches(2.9)
gap = Inches(0.25)
start_x = Inches(0.6)
top = Inches(1.7)
box_h = Inches(4.6)
for i, (title, body) in enumerate(stages):
    x = start_x + i * (box_w + gap)
    add_rect(s, x, top, box_w, box_h, LGREEN if i % 2 == 0 else LGRAY)
    add_rect(s, x, top, box_w, Inches(0.6), TEAL)
    add_text(s, x, top, box_w, Inches(0.6), title, size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), top + Inches(0.75), box_w - Inches(0.3), box_h - Inches(0.9),
             body, size=12.5, color=DGRAY)
    if i < len(stages) - 1:
        ax = x + box_w + Emu(0)
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax + Inches(0.02), top + Inches(2.0),
                                     gap - Inches(0.04), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background()
        arrow.shadow.inherit = False

add_text(s, Inches(0.6), Inches(6.45), Inches(12.0), Inches(0.7),
         "Fully offline by default (TF-IDF + extractive). Claude path only used if "
         "ANTHROPIC_API_KEY is configured — same retrieval layer, more fluent answers.",
         size=13, color=DGRAY, italic=True)
add_footer(s, 14)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 15 — 15-QUERY EVALUATION
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Evaluation — 15 Representative Queries", track="rag")

# Donut-ish summary as a simple bar/table
rows = len(EVAL_SUMMARY) + 1
tbl = s.shapes.add_table(rows, 3, Inches(0.6), Inches(1.5), Inches(5.6), Inches(2.6)).table
for c, txt in enumerate(["Label", "Count", "Share"]):
    cell = tbl.cell(0, c)
    cell.text = txt
    cell.fill.solid()
    cell.fill.fore_color.rgb = TEAL
    cell.text_frame.paragraphs[0].runs[0].font.bold = True
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
for r, (label, count, share) in enumerate(EVAL_SUMMARY, start=1):
    for c, txt in enumerate([label, str(count), share]):
        cell = tbl.cell(r, c)
        cell.text = txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = LGRAY if r % 2 == 0 else WHITE
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(14)
        run.font.color.rgb = DGRAY
        run.font.bold = (c == 0)

add_bullets(s, Inches(0.6), Inches(4.3), Inches(5.6), Inches(2.9), [
    "Correct (12/15): single-document queries — SOPs, ISO mapping, model "
    "performance, maintenance, deployment — all retrieved and answered correctly",
    "Partial (2/15): severity comparison & confidence threshold — info split "
    "across two documents, not always combined",
    "No answer (1/15): out-of-scope query (cotton yarn price) — system correctly "
    "returns low-relevance docs, doesn't fabricate an answer",
    "Hallucinated (0/15): extractive composer can only return real corpus "
    "sentences; Claude path is grounded by “context-only” prompt",
], size=13, space_after=8)

add_text(s, Inches(6.6), Inches(1.5), Inches(6.2), Inches(0.5),
         "Sample queries", size=16, bold=True, color=TEAL)
add_bullets(s, Inches(6.6), Inches(2.05), Inches(6.2), Inches(5.0), [
    "“What causes a Circle defect and how do I fix it?” → Correct",
    "“Why is this a data-driven Digital Twin rather than a simulation?” → Correct",
    "“Why does per-fold accuracy vary between 83% and 100%?” → Correct",
    "“How would this integrate with a factory MES/QMS?” → Correct",
    "“Is a Circle defect more severe than a Line defect?” → Partial",
    "“What confidence threshold should trigger human review?” → Partial",
    "“What is the current market price of cotton yarn?” → No answer",
], size=14, space_after=10)
add_footer(s, 15)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 16 — BENEFITS / LIMITATIONS / DEPLOYMENT CHALLENGES
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Benefits, Limitations and Deployment Challenges", track="rag")

col_w = Inches(4.0)
gap = Inches(0.15)
cols_data = [
    ("Benefits", GREEN, [
        "Closes the loop: defect prediction → relevant SOP, no manual lookup",
        "Fully offline by default — no API cost, no data leaves the factory",
        "Source-cited answers with relevance scores — verifiable",
        "Optional Claude upgrade without changing retrieval/corpus",
    ]),
    ("Limitations", ORANGE, [
        "12-doc corpus is illustrative, not sourced from real mill manuals",
        "TF-IDF is lexical — can miss paraphrased queries",
        "Weak multi-document synthesis (seen in “partial” results)",
        "No feedback loop from operator corrections into the corpus",
    ]),
    ("Deployment Challenges", TEAL, [
        "Corpus maintenance — keeping SOPs/standards up to date",
        "Access control — some SOPs may be commercially sensitive",
        "Integration with Cross System Entity (live MES/QMS records)",
        "Offline fallback needed for unreliable shop-floor connectivity",
    ]),
]
for i, (title, color, items) in enumerate(cols_data):
    x = Inches(0.5) + i * (col_w + gap)
    add_rect(s, x, Inches(1.4), col_w, Inches(0.55), color)
    add_text(s, x, Inches(1.4), col_w, Inches(0.55), title, size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, x + Inches(0.1), Inches(2.1), col_w - Inches(0.2), Inches(5.0),
                items, size=13, space_after=10)
add_footer(s, 16)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 17 — CONCLUSIONS
# ════════════════════════════════════════════════════════════════════════
s = content_slide("Conclusions")
add_bullets(s, Inches(0.6), Inches(1.4), Inches(12.0), Inches(5.6), [
    f"Track 2: ResNet18 (frozen backbone, 3-class head) reaches "
    f"{RES['avg_accuracy']:.1%} (±{RES['std_accuracy']:.1%}) average accuracy "
    "on stratified 5-fold CV over 150 images, with Grad-CAM giving the operator "
    "a per-prediction view of model attention",
    "Track 3: 12-document SOP/standards knowledge base via TF-IDF retrieval, "
    "achieved 12/15 correct, 2/15 partial, 1/15 no-answer, 0/15 hallucinated on "
    "the required 15-query evaluation",
    "Both tracks integrated in a single 5-page Streamlit dashboard — a defect "
    "prediction automatically triggers the relevant SOP lookup (detection-to-action)",
    "ISO 23247 mapping situates both components as parts of one Digital Twin: "
    "OME, DCE, Digital Twin Entity, User Entity, Cross System Entity",
    "Honest limitations documented: small/staged dataset, illustrative SOP corpus, "
    "single fabric pattern — with concrete recommendations for moving from "
    "prototype to pilot (confidence-threshold human review, retraining pipeline, "
    "real mill documentation)",
], size=17, space_after=16)
add_footer(s, 17)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 18 — DEMO / THANK YOU
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, DGREEN)
add_text(s, Inches(1), Inches(2.6), Inches(11.33), Inches(1.0),
         "Live Demo", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.7), Inches(11.33), Inches(0.8),
         "streamlit run app.py", size=22, color=RGBColor(0xB8, 0xE0, 0xC8),
         align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, Inches(1), Inches(4.6), Inches(11.33), Inches(0.6),
         "EDA  ·  Defect Classifier + Grad-CAM  ·  Model Performance  ·  Quality Assistant",
         size=15, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(1), Inches(6.0), Inches(11.33), Inches(0.6),
         "Thank you", size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


prs.save(OUT_PPTX)
print(f"Presentation written to: {OUT_PPTX}")
